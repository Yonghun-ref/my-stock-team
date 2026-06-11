"""report-pptx 렌더러: 구조화 JSON → KB 디자인 PPTX.

사용법:
    python render.py <input.json> <output.pptx>

설계 원칙(스킬 사양):
- 슬라이드 순서 고정: 표지 → 종목 개요 → 재무 요약 → 가격/추세 → 뉴스·심리 → 리스크 → 한 줄 종합
- 포인트색 KB 옐로우(#FFBC00) + 본문 그레이/화이트, 차분한 금융 리포트 톤
- 한글 폰트는 '맑은 고딕' 하나로 고정(latin/ea/cs 모두 지정 → 글자 깨짐 방지)
- 표가 슬라이드 밖으로 넘치면 행을 줄이고 생략 표기
- 출처·매수/매도 가드레일은 입력 JSON 단계에서 보장(여기선 받은 값만 렌더)
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── 디자인 토큰 ─────────────────────────────────────────────
FONT = "맑은 고딕"
KB_YELLOW = RGBColor(0xFF, 0xBC, 0x00)
INK = RGBColor(0x2B, 0x2B, 0x2B)        # 본문 진회색
SUB = RGBColor(0x70, 0x70, 0x70)        # 보조 회색
LINE = RGBColor(0xD9, 0xD9, 0xD9)       # 옅은 경계선
HDR_BG = RGBColor(0x33, 0x33, 0x33)     # 표 헤더 배경(다크그레이)
ROW_ALT = RGBColor(0xF5, 0xF5, 0xF5)    # 표 줄무늬
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DISCLAIMER = "본 리포트는 학습용 분석 자료이며, 투자 권유나 자문이 아닙니다."

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.7)
BODY_W = SLIDE_W - MARGIN * 2

# 표 넘침 방지: 본문 영역 높이/행높이로 최대 행 수 산정
MAX_TABLE_ROWS = 9  # 헤더 포함. 초과 시 잘라내고 생략 표기


def set_font(run, *, size=None, bold=None, color=None, name=FONT):
    """latin/ea/cs typeface를 모두 지정해 한글 깨짐을 막는다."""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _rect(slide, left, top, width, height, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    return shp


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _section_header(slide, label, title):
    """상단 섹션 라벨 + 제목 + KB 옐로우 룰."""
    _, tf = _textbox(slide, MARGIN, Inches(0.45), BODY_W, Inches(0.4))
    r = tf.paragraphs[0].add_run()
    r.text = label
    set_font(r, size=12, bold=True, color=KB_YELLOW)

    _, tf2 = _textbox(slide, MARGIN, Inches(0.78), BODY_W, Inches(0.7))
    r2 = tf2.paragraphs[0].add_run()
    r2.text = title
    set_font(r2, size=26, bold=True, color=INK)

    _rect(slide, MARGIN, Inches(1.5), BODY_W, Pt(3), fill=KB_YELLOW)


def _footer(slide, page):
    _, tf = _textbox(slide, MARGIN, Inches(7.02), BODY_W, Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = DISCLAIMER
    set_font(r, size=9, color=SUB)
    # 페이지 번호
    _, tf2 = _textbox(slide, SLIDE_W - Inches(1.2), Inches(7.02), Inches(0.6), Inches(0.35))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(page)
    set_font(r2, size=9, color=SUB)


def add_bullets(slide, top, items, *, size=14, max_items=8):
    if not items:
        return top
    items = list(items)
    extra = 0
    if len(items) > max_items:
        extra = len(items) - max_items
        items = items[:max_items]
    height = Inches(0.42) * (len(items) + (1 if extra else 0)) + Inches(0.2)
    _, tf = _textbox(slide, MARGIN, top, BODY_W, height)
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        dot = p.add_run()
        dot.text = "▪ "
        set_font(dot, size=size, color=KB_YELLOW, bold=True)
        r = p.add_run()
        r.text = str(text)
        set_font(r, size=size, color=INK)
    if extra:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"… 외 {extra}개 항목 (전체는 리서치 .md 참조)"
        set_font(r, size=11, color=SUB)
    return top + height


def add_table(slide, top, table, *, max_rows=MAX_TABLE_ROWS):
    """markdown 표(dict: headers, rows) → 디자인 표. 넘치면 행 축소."""
    if not table or not table.get("headers"):
        return top
    headers = table["headers"]
    rows = [list(map(str, r)) for r in table.get("rows", [])]

    truncated = 0
    if len(rows) > max_rows - 1:
        truncated = len(rows) - (max_rows - 1)
        rows = rows[: max_rows - 1]

    ncols = len(headers)
    nrows = len(rows) + 1
    row_h = Inches(0.4)
    tbl_h = row_h * nrows
    gobj = slide.shapes.add_table(nrows, ncols, MARGIN, top, BODY_W, tbl_h)
    tbl = gobj.table

    # 컬럼 폭: 첫 컬럼 넓게, 나머지 균등
    if ncols > 1:
        first = int(BODY_W * 0.30)
        rest = int((BODY_W - first) / (ncols - 1))
        tbl.columns[0].width = first
        for c in range(1, ncols):
            tbl.columns[c].width = rest

    # 컬럼 수가 많으면 폰트 축소(넘침 방지)
    body_size = 12 if ncols <= 4 else (10 if ncols <= 6 else 9)
    hdr_size = body_size + 1

    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HDR_BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)
        para = cell.text_frame.paragraphs[0]
        r = para.add_run()
        r.text = str(h)
        set_font(r, size=hdr_size, bold=True, color=WHITE)

    for ri, row in enumerate(rows, start=1):
        for c in range(ncols):
            cell = tbl.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if ri % 2 == 0 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)
            para = cell.text_frame.paragraphs[0]
            r = para.add_run()
            r.text = row[c] if c < len(row) else ""
            set_font(r, size=body_size, color=INK)

    bottom = top + tbl_h
    if truncated:
        _, tf = _textbox(slide, MARGIN, bottom + Inches(0.05), BODY_W, Inches(0.3))
        r = tf.paragraphs[0].add_run()
        r.text = f"… 외 {truncated}개 행 생략 (전체는 리서치 .md 참조)"
        set_font(r, size=10, color=SUB)
        bottom += Inches(0.35)
    return bottom + Inches(0.15)


def add_notes(slide, top, notes, *, size=11):
    if not notes:
        return top
    _, tf = _textbox(slide, MARGIN, top, BODY_W, Inches(0.4) * len(notes) + Inches(0.1))
    for i, n in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"※ {n}"
        set_font(r, size=size, color=SUB)
    return top


# ── 슬라이드 빌더 ───────────────────────────────────────────
def build(data, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    name = data.get("ticker_name", "종목 미상")
    date = data.get("date", "")
    opinion = data.get("opinion", "")

    # 1) 표지
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, SLIDE_W, Inches(2.6), fill=KB_YELLOW)
    _rect(s, 0, Inches(2.6), SLIDE_W, Pt(4), fill=INK)
    _, tf = _textbox(s, MARGIN, Inches(0.85), BODY_W, Inches(1.0))
    r = tf.paragraphs[0].add_run()
    r.text = name
    set_font(r, size=40, bold=True, color=INK)
    _, tf2 = _textbox(s, MARGIN, Inches(1.85), BODY_W, Inches(0.5))
    r2 = tf2.paragraphs[0].add_run()
    r2.text = f"분석 리서치  ·  작성일 {date}" if date else "분석 리서치"
    set_font(r2, size=16, color=INK)
    if opinion:
        _, tf3 = _textbox(s, MARGIN, Inches(3.2), BODY_W, Inches(0.6))
        lab = tf3.paragraphs[0].add_run()
        lab.text = "종합의견  "
        set_font(lab, size=18, color=SUB)
        val = tf3.paragraphs[0].add_run()
        val.text = opinion
        set_font(val, size=18, bold=True, color=INK)
    _, tf4 = _textbox(s, MARGIN, Inches(6.7), BODY_W, Inches(0.5))
    r4 = tf4.paragraphs[0].add_run()
    r4.text = DISCLAIMER
    set_font(r4, size=10, color=SUB)

    page = 1

    def content_slide(label, title):
        nonlocal page
        page += 1
        sl = prs.slides.add_slide(blank)
        _section_header(sl, label, title)
        _footer(sl, page)
        return sl

    top0 = Inches(1.75)

    # 2) 종목 개요
    ov = data.get("overview", {})
    sl = content_slide("OVERVIEW", "종목 개요")
    top = add_bullets(sl, top0, ov.get("bullets"))
    top = add_table(sl, top, ov.get("table"))
    add_notes(sl, top, ov.get("notes"))

    # 3) 재무 요약
    fin = data.get("financials", {})
    sl = content_slide("FINANCIALS", "재무 요약 (최근 3개년)")
    top = add_table(sl, top0, fin.get("table"))
    top = add_bullets(sl, top, fin.get("bullets"), size=13)
    add_notes(sl, top, fin.get("notes"))

    # 4) 가격/추세
    pr = data.get("price", {})
    sl = content_slide("PRICE & TREND", "가격 / 추세")
    top = add_table(sl, top0, pr.get("table"))
    top = add_bullets(sl, top, pr.get("bullets"), size=13)
    add_notes(sl, top, pr.get("notes"))

    # 5) 뉴스·심리
    nw = data.get("news", {})
    sl = content_slide("NEWS & SENTIMENT", "뉴스 · 심리")
    top = add_bullets(sl, top0, nw.get("bullets"))
    if nw.get("sentiment"):
        _rect(sl, MARGIN, top + Inches(0.1), BODY_W, Inches(0.7), fill=ROW_ALT, line=LINE)
        _, tf = _textbox(sl, MARGIN + Inches(0.2), top + Inches(0.22), BODY_W - Inches(0.4), Inches(0.5))
        lab = tf.paragraphs[0].add_run()
        lab.text = "시장 심리  "
        set_font(lab, size=13, bold=True, color=KB_YELLOW)
        val = tf.paragraphs[0].add_run()
        val.text = nw["sentiment"]
        set_font(val, size=13, color=INK)

    # 6) 리스크
    rk = data.get("risk", {})
    sl = content_slide("RISK", "리스크")
    top = add_bullets(sl, top0, rk.get("bullets"), max_items=7)
    add_notes(sl, top, rk.get("notes"))

    # 7) 한 줄 종합
    sm = data.get("summary", {})
    sl = content_slide("CONCLUSION", "한 줄 종합")
    _rect(sl, MARGIN, top0, Pt(5), Inches(1.1), fill=KB_YELLOW)
    _, tf = _textbox(sl, MARGIN + Inches(0.25), top0, BODY_W - Inches(0.25), Inches(1.2))
    r = tf.paragraphs[0].add_run()
    r.text = sm.get("headline", "")
    set_font(r, size=20, bold=True, color=INK)
    add_bullets(sl, top0 + Inches(1.4), sm.get("bullets"), size=14)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return len(prs.slides._sldIdLst)


def main():
    if len(sys.argv) != 3:
        print("usage: python render.py <input.json> <output.pptx>", file=sys.stderr)
        sys.exit(2)
    data = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))
    n = build(data, sys.argv[2])
    print(f"OK slides={n} -> {sys.argv[2]}")


if __name__ == "__main__":
    main()
